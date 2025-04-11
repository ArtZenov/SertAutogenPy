import pandas as pd
from pptx import Presentation
import os
from pptx.util import Pt
import time
import comtypes.client
import argparse
import logging
import yagmail
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import Mail, Attachment, FileContent, FileName, FileType, Disposition
import base64
from config import SENDGRID_API_KEY, GMAIL_USER, GMAIL_PASSWORD

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
            logging.FileHandler("cert_log.txt", encoding="utf-8"),
            logging.StreamHandler()
        ]
)

def fit_text_to_shape(shape, text, max_width=None):
    shape.text = text
    if not shape.text_frame.paragraphs:
        return shape
    font_size = 40
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    for run in paragraph.runs:
        run.font.size = Pt(font_size)
    if max_width is None:
        max_width = shape.width
    char_width_factor = 0.5
    while (len(text) * Pt(font_size).inches * char_width_factor > max_width.inches) and font_size > 8:
        font_size -= 2
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
    shape.text_frame.word_wrap = False
    return shape

def convert_pptx_to_pdf(pptx_path, pdf_path):
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
    presentation.SaveAs(os.path.abspath(pdf_path), 32)
    presentation.Close()
    powerpoint.Quit()

def generate_certificates(data, template_ppt, limit=None):
    if limit:
        data = data.iloc[:limit]
    os.makedirs("certificates", exist_ok=True)
    for index, row in data.iterrows():
        name = row['ФИО']
        email = row['Email']
        if pd.isna(name) or pd.isna(email):
            logging.warning(f"Пропущены данные для строки {index + 1}: ФИО={name}, Email={email}")
            continue
        prs = Presentation(template_ppt)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and "{ФИО}" in shape.text:
                    new_text = shape.text.replace("{ФИО}", name)
                    shape = fit_text_to_shape(shape, new_text)
        cert_name = f"Сертификат_{name.replace(' ', '_')}".replace('/', '_').replace('\\', '_')
        pptx_path = os.path.join("certificates", f"{cert_name}.pptx")
        pdf_path = os.path.join("certificates", f"{cert_name}.pdf")
        try:
            prs.save(pptx_path)
            convert_pptx_to_pdf(pptx_path, pdf_path)
            os.remove(pptx_path)
            logging.info(f"Сертификат создан: {pdf_path}")
        except Exception as e:
            logging.error(f"Ошибка при создании сертификата для {name}: {e}")
            continue

def send_certificates(data, limit=None, email_service='gmail', start_index=0):
    if limit:
        data = data.iloc[start_index:start_index + limit]
    else:
        data = data.iloc[start_index:]

    sent_emails = set()
    try:
        with open('cert_log.txt', 'r', encoding='utf-8') as log_file:
            for line in log_file:
                if "Сертификат отправлен" in line:
                    email = line.split('(')[-1].strip(')\n')
                    sent_emails.add(email)
    except FileNotFoundError:
        pass

    if email_service == 'sendgrid':
        sg = SendGridAPIClient(SENDGRID_API_KEY)
    elif email_service == 'gmail':
        yag = yagmail.SMTP(GMAIL_USER, GMAIL_PASSWORD)

    for index, row in data.iterrows():
        name = row['ФИО']
        email = row['Email']
        if pd.isna(name) or pd.isna(email):
            logging.warning(f"Пропущены данные для строки {index + 1}: ФИО={name}, Email={email}")
            continue
        if email in sent_emails:
            logging.info(f"Сертификат уже отправлен ранее: {name} ({email})")
            continue
        cert_name = f"Сертификат_{name.replace(' ', '_')}".replace('/', '_').replace('\\', '_')
        pdf_path = os.path.join("certificates", f"{cert_name}.pdf")
        if not os.path.exists(pdf_path):
            logging.warning(f"PDF не найден для {name}: {pdf_path}")
            continue

        try:
            if email_service == 'sendgrid':
                message = Mail(
                    from_email=GMAIL_USER,
                    to_emails=email,
                    subject="Ваш сертификат участника Лингва.Тех",
                    html_content=f"Здравствуйте, {name}!<br>Во вложении — ваш сертификат участника."
                )
                with open(pdf_path, 'rb') as f:
                    data = f.read()
                encoded_file = base64.b64encode(data).decode()
                attachment = Attachment(
                    FileContent(encoded_file),
                    FileName(os.path.basename(pdf_path)),
                    FileType('application/pdf'),
                    Disposition('attachment')
                )
                message.attachment = attachment
                sg.send(message)
            elif email_service == 'gmail':
                yag.send(
                    to=email,
                    subject="Ваш сертификат участника Лингва.Тех",
                    contents=f"Здравствуйте, {name}!\n\nВо вложении — ваш сертификат участника.",
                    attachments=pdf_path
                )
            logging.info(f"Сертификат отправлен: {name} ({email})")
        except Exception as e:
            logging.error(f"Ошибка при отправке письма для {name}: {e}")
        time.sleep(5)

    if email_service == 'gmail':
        yag.close()

def main():
    parser = argparse.ArgumentParser(description="Генерация и рассылка сертификатов")
    parser.add_argument('--limit', type=int, default=None, help="Ограничение на количество сертификатов")
    parser.add_argument('--mode', choices=['all', 'generate', 'send'], default='all',
                        help="Режим работы: all (всё), generate (только создание), send (только отправка)")
    parser.add_argument('--email_service', choices=['gmail', 'sendgrid'], default='gmail',
                        help="Сервис отправки: gmail или sendgrid")
    parser.add_argument('--start_index', type=int, default=0,
                        help="Индекс строки, с которой начать отправку (для send)")
    args = parser.parse_args()

    try:
        data = pd.read_excel("data/list.xlsx")
    except FileNotFoundError:
        logging.error("Файл test_list.xlsx не найден.")
        exit(1)

    if not {'ФИО', 'Email'}.issubset(data.columns):
        logging.error("В файле test_list.xlsx отсутствуют столбцы 'ФИО' или 'Email'.")
        exit(1)

    template_ppt = "data/sert.pptx"
    try:
        Presentation(template_ppt)
    except FileNotFoundError:
        logging.error("Файл sert.pptx не найден.")
        exit(1)
    except Exception as e:
        logging.error(f"Ошибка при загрузке шаблона: {e}")
        exit(1)

    if args.mode in ['all', 'generate']:
        logging.info("Начало генерации сертификатов")
        generate_certificates(data, template_ppt, args.limit)
    if args.mode in ['all', 'send']:
        logging.info("Начало отправки сертификатов")
        send_certificates(data, args.limit, args.email_service, args.start_index)

    logging.info("Работа завершена")

if __name__ == "__main__":
    main()